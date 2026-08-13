"""Tests for the scientific-slides tooling.

`generate_schematic.py` and `generate_schematic_ai.py` are byte-identical to the
copies in four other skills, so their behaviour comes from the shared contract
(`skill_contract.schematic`) rather than being re-tested here. What is specific
to this skill is the slide pipeline, and each group of tests below guards one way
it fails quietly:

* `validate_presentation.py` is the gate before a talk is delivered. Its slide
  budget is checked in both directions -- a deck inside the range must be
  reported as fine, one outside it must be warned about -- and against the
  one-slide-per-minute rule its own `assets/timing_guidelines.md` documents. The
  PDF path is driven with real two-page PDFs whose 16:9 and 4:3 geometry is known
  by construction, which is also what caught a bug: the script imported the
  retired `PyPDF2` only, so with the maintained `pypdf` installed every PDF check
  silently degraded to "not installed".
* `pdf_to_images.py` converts at a requested DPI. 144 DPI on a 720-point page is
  1440 pixels, so the zoom arithmetic is checked against the rendered size rather
  than assumed.
* `slides_to_pdf.py` decides slide *order* from filenames, so `02` must land
  before `10`; a lexicographic sort would silently reorder a talk.
* The AI scripts are exercised without touching the network: the OpenRouter
  request is mocked, and what is asserted is the quality-threshold decision, the
  response parsing, the retry loop, and that the API key travels in the child's
  environment rather than on its command line.
"""

from __future__ import annotations

import base64
import contextlib
import io
import os
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path
from unittest import mock

import pytest

import skill_contract

SKILL_ROOT = Path(__file__).resolve().parents[2] / "skills" / "scientific-slides"
SCRIPTS = SKILL_ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

# `slides_to_pdf` exits at import time when Pillow is missing, so the guard has
# to come first: without it a bare environment turns a skip into a collection
# error for the whole module.
Image = pytest.importorskip("PIL.Image", reason="the slide scripts need Pillow")

import generate_slide_image  # noqa: E402  (standard library only)
import pdf_to_images  # noqa: E402
import slides_to_pdf  # noqa: E402
import validate_presentation  # noqa: E402

SchematicTests = skill_contract.schematic.schematic_test_case(SKILL_ROOT)
ReviewParsingTests = skill_contract.schematic.review_parsing_test_case(
    SCRIPTS, "generate_schematic_ai"
)
ReviewFailureTests = skill_contract.schematic.review_failure_test_case(
    SCRIPTS, "generate_schematic_ai", "ScientificSchematicGenerator",
    ("diagram.png", "a prompt", 1, "journal", 2),
)

# `generate_slide_image_ai.py` is a fork of the schematic generator, not one of
# the byte-identical copies, so it needs its own instantiation of the same
# review-parsing contract.
SlideImageReviewParsingTests = skill_contract.schematic.review_parsing_test_case(
    SCRIPTS, "generate_slide_image_ai"
)
SlideImageReviewFailureTests = skill_contract.schematic.review_failure_test_case(
    SCRIPTS, "generate_slide_image_ai", "SlideImageGenerator",
    ("slide.png", "a prompt", 1, False, 2),
)
CliHelpTests = skill_contract.cli.help_test_case(SKILL_ROOT)


def child_environment() -> dict[str, str]:
    """This process's environment minus any real OpenRouter credential."""
    environment = {**os.environ, "PYTHONDONTWRITEBYTECODE": "1"}
    environment.pop("OPENROUTER_API_KEY", None)
    return environment


class TemporaryDirectoryTestCase(unittest.TestCase):
    def setUp(self) -> None:
        self._temporary = tempfile.TemporaryDirectory()
        self.addCleanup(self._temporary.cleanup)
        self.root = Path(self._temporary.name)

    def deck(self, pages: int = 2, size: tuple[int, int] = (1280, 720)) -> Path:
        """A PDF of `pages` blank slides, one point per pixel at 72 dpi."""
        images = [Image.new("RGB", size, (255, 255, 255)) for _ in range(pages)]
        path = self.root / f"deck_{pages}_{size[0]}x{size[1]}.pdf"
        images[0].save(
            path, "PDF", resolution=72, save_all=True, append_images=images[1:]
        )
        return path

    #: Fill colour per image mode -- Pillow wants a scalar for single-band modes.
    FILL = {"RGB": (200, 200, 200), "RGBA": (200, 200, 200, 255), "L": 200, "P": 200}

    def image(self, name: str, size: tuple[int, int] = (320, 180), mode="RGB") -> Path:
        path = self.root / name
        path.parent.mkdir(parents=True, exist_ok=True)
        Image.new(mode, size, self.FILL[mode]).save(path)
        return path


def validated(path, duration=None) -> dict:
    validator = validate_presentation.PresentationValidator(str(path), duration)
    with contextlib.redirect_stdout(io.StringIO()):
        return validator.validate()


class SlideBudgetTests(unittest.TestCase):
    GUIDELINES = validate_presentation.PresentationValidator.SLIDE_GUIDELINES

    def counted(self, slides: int, duration: int) -> dict:
        validator = validate_presentation.PresentationValidator("talk.pdf", duration)
        validator._check_slide_count(slides)
        return {"info": validator.info, "warnings": validator.warnings}

    def test_every_duration_brackets_its_recommendation(self) -> None:
        for duration, (minimum, recommended, maximum) in self.GUIDELINES.items():
            with self.subTest(duration=duration):
                self.assertLess(minimum, recommended)
                self.assertLess(recommended, maximum)

    def test_the_budget_grows_with_the_talk_length(self) -> None:
        recommended = [self.GUIDELINES[key][1] for key in sorted(self.GUIDELINES)]
        self.assertEqual(recommended, sorted(recommended))

    def test_the_recommendations_track_one_slide_per_minute(self) -> None:
        # The rule the skill's own timing_guidelines.md is built on; a table that
        # drifted far from it would contradict the guidance it ships with.
        for duration, (_, recommended, _) in self.GUIDELINES.items():
            with self.subTest(duration=duration):
                self.assertGreaterEqual(recommended / duration, 0.8)
                self.assertLessEqual(recommended / duration, 1.25)

    def test_the_ranges_overlap_the_shipped_timing_asset(self) -> None:
        # assets/timing_guidelines.md tabulates "Duration | Total Slides"; the
        # validator must not call a deck too long when that table calls it right.
        table = (SKILL_ROOT / "assets" / "timing_guidelines.md").read_text(
            encoding="utf-8"
        )
        import re

        rows = re.findall(r"^\|\s*(\d+)\s*min\s*\|\s*(\d+)-(\d+)", table, re.MULTILINE)
        self.assertTrue(rows, "the timing asset no longer tabulates slide counts")
        for duration, low, high in rows:
            duration, low, high = int(duration), int(low), int(high)
            if duration not in self.GUIDELINES:
                continue
            minimum, _, maximum = self.GUIDELINES[duration]
            with self.subTest(duration=duration):
                self.assertLessEqual(minimum, high)
                self.assertGreaterEqual(maximum, low)

    def test_a_deck_inside_the_range_is_reported_as_fine(self) -> None:
        # 16 slides for a 15-minute talk: squarely inside 13-20.
        result = self.counted(16, 15)
        self.assertEqual(result["warnings"], [])
        self.assertTrue(
            any("within recommended range" in line for line in result["info"])
        )

    def test_too_few_slides_warns_about_the_empty_time(self) -> None:
        result = self.counted(4, 15)
        self.assertTrue(result["warnings"])
        self.assertIn("Fewer slides", result["warnings"][0])

    def test_too_many_slides_warns_about_running_over(self) -> None:
        result = self.counted(40, 15)
        self.assertTrue(result["warnings"])
        self.assertIn("run over time", result["warnings"][0])

    def test_the_boundaries_themselves_are_acceptable(self) -> None:
        minimum, _, maximum = self.GUIDELINES[15]
        for slides in (minimum, maximum):
            with self.subTest(slides=slides):
                self.assertEqual(self.counted(slides, 15)["warnings"], [])

    def test_an_untabulated_duration_borrows_the_nearest_one(self) -> None:
        # 12 minutes is closest to the 10-minute row, and the substitution has
        # to be stated or the numbers look wrong to the user.
        result = self.counted(11, 12)
        self.assertTrue(
            any("guidelines for 10-minute" in line for line in result["info"]),
            result["info"],
        )


class ValidatorFileTests(TemporaryDirectoryTestCase):
    def test_a_missing_file_is_the_only_thing_reported(self) -> None:
        result = validated(self.root / "absent.pdf")
        self.assertFalse(result["valid"])
        self.assertIn("File not found", result["issues"][0])

    def test_an_unsupported_extension_warns_but_does_not_fail(self) -> None:
        path = self.root / "talk.key"
        path.write_bytes(b"")
        result = validated(path)
        self.assertTrue(result["valid"])
        self.assertTrue(any("Unknown file type" in line for line in result["warnings"]))

    def test_the_file_type_is_recognised_case_insensitively(self) -> None:
        validator = validate_presentation.PresentationValidator("TALK.PDF")
        self.assertEqual(validator.file_type, ".pdf")

    def sized(self, megabytes: float) -> dict:
        path = self.root / f"deck_{megabytes}.pdf"
        with path.open("wb") as handle:
            handle.truncate(int(megabytes * 1024 * 1024))  # sparse, not written
        validator = validate_presentation.PresentationValidator(str(path))
        validator._check_file_size()
        return {"info": validator.info, "warnings": validator.warnings,
                "issues": validator.issues}

    def test_an_ordinary_deck_size_is_only_reported(self) -> None:
        result = self.sized(4)
        self.assertEqual(result["warnings"], [])
        self.assertEqual(result["issues"], [])
        self.assertIn("4.00 MB", result["info"][0])

    def test_a_large_deck_warns_about_sharing_it(self) -> None:
        result = self.sized(60)
        self.assertTrue(result["warnings"])
        self.assertEqual(result["issues"], [])

    def test_an_enormous_deck_is_an_issue_not_a_warning(self) -> None:
        result = self.sized(120)
        self.assertTrue(result["issues"])
        self.assertIn("compressing images", result["issues"][0])


class ValidatorPdfTests(TemporaryDirectoryTestCase):
    def setUp(self) -> None:
        super().setUp()
        if not validate_presentation.HAS_PDF_READER:
            self.skipTest("neither pypdf nor PyPDF2 is installed")

    def test_a_widescreen_deck_reports_its_page_count_and_ratio(self) -> None:
        # Regression: with only pypdf installed this whole branch used to be
        # skipped, and the validator reported nothing about the deck at all.
        result = validated(self.deck(pages=2, size=(1280, 720)))
        self.assertTrue(result["valid"], result["issues"])
        joined = "\n".join(result["info"])
        self.assertIn("Number of slides: 2", joined)
        self.assertIn("16:9", joined)

    def test_a_four_by_three_deck_is_recognised(self) -> None:
        result = validated(self.deck(pages=1, size=(1024, 768)))
        self.assertIn("4:3", "\n".join(result["info"]))

    def test_an_unusual_aspect_ratio_is_flagged_for_the_venue(self) -> None:
        result = validated(self.deck(pages=1, size=(1000, 1000)))
        self.assertTrue(
            any("Unusual aspect ratio" in line for line in result["warnings"])
        )

    def test_the_slide_dimensions_are_reported_in_inches(self) -> None:
        # 1280 x 720 points at 72 points per inch is 17.8" x 10.0".
        result = validated(self.deck(pages=1, size=(1280, 720)))
        self.assertIn('17.8" × 10.0"', "\n".join(result["info"]))

    def test_the_page_count_is_measured_against_the_duration(self) -> None:
        result = validated(self.deck(pages=2, size=(1280, 720)), duration=15)
        self.assertTrue(any("Fewer slides" in line for line in result["warnings"]))

    def test_a_corrupt_pdf_is_an_issue_rather_than_a_traceback(self) -> None:
        path = self.root / "broken.pdf"
        path.write_bytes(b"%PDF-1.4\nnot really a pdf\n")
        result = validated(path)
        self.assertFalse(result["valid"])
        self.assertIn("Error reading PDF", result["issues"][0])


class ValidatorLatexTests(TemporaryDirectoryTestCase):
    """The compile step is stubbed: whether pdflatex exists is not the subject."""

    def source(self) -> Path:
        path = self.root / "talk.tex"
        path.write_text("\\documentclass{beamer}\\begin{document}\\end{document}")
        return path

    def test_a_failed_compile_is_an_issue_pointing_at_the_log(self) -> None:
        path = self.source()
        with mock.patch("subprocess.run", side_effect=FileNotFoundError):
            result = validated(path)
        self.assertFalse(result["valid"])
        self.assertIn(".log", result["issues"][0])

    def test_a_successful_compile_folds_in_the_pdf_report(self) -> None:
        if not validate_presentation.HAS_PDF_READER:
            self.skipTest("neither pypdf nor PyPDF2 is installed")
        path = self.source()
        # pdflatex writes talk.pdf beside talk.tex; produce one it can find.
        self.deck(pages=3, size=(1280, 720)).replace(path.with_suffix(".pdf"))
        completed = subprocess.CompletedProcess(args=[], returncode=0)
        with mock.patch("subprocess.run", return_value=completed):
            result = validated(path)
        joined = "\n".join(result["info"])
        self.assertIn("LaTeX compilation: SUCCESS", joined)
        # The merged PDF report is what makes the .tex path useful.
        self.assertIn("Number of slides: 3", joined)

    def test_the_compile_never_enables_shell_escape(self) -> None:
        # Shell escape in a LaTeX source would run arbitrary commands.
        path = self.source()
        with mock.patch("subprocess.run", side_effect=FileNotFoundError) as run:
            validated(path)
        self.assertIn("-no-shell-escape", run.call_args.args[0])
        self.assertIn("-interaction=nonstopmode", run.call_args.args[0])


class ValidatorPptxTests(TemporaryDirectoryTestCase):
    def setUp(self) -> None:
        super().setUp()
        if not validate_presentation.HAS_PPTX:
            self.skipTest("python-pptx is not installed")

    def build(self, *, font_pt: int, bullets: int) -> Path:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        frame = box.text_frame
        for index in range(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            run = paragraph.add_run()
            run.text = f"point {index}"
            run.font.size = Pt(font_pt)
        path = self.root / f"deck_{font_pt}_{bullets}.pptx"
        presentation.save(path)
        return path

    def test_a_readable_uncluttered_slide_raises_nothing(self) -> None:
        result = validated(self.build(font_pt=24, bullets=3))
        self.assertTrue(result["valid"])
        self.assertEqual(
            [line for line in result["warnings"] if "text" in line or "bullets" in line],
            [],
        )

    def test_text_below_eighteen_point_is_flagged_with_the_slide_number(self) -> None:
        # 18 pt is the projection floor the skill documents.
        result = validated(self.build(font_pt=12, bullets=2))
        self.assertTrue(
            any("Small text (<18pt)" in line and "[1]" in line
                for line in result["warnings"]),
            result["warnings"],
        )

    def test_eighteen_point_itself_is_acceptable(self) -> None:
        result = validated(self.build(font_pt=18, bullets=2))
        self.assertEqual(
            [line for line in result["warnings"] if "Small text" in line], []
        )

    def test_more_than_six_bullets_is_flagged(self) -> None:
        result = validated(self.build(font_pt=24, bullets=8))
        self.assertTrue(
            any("Many bullets (>6)" in line for line in result["warnings"]),
            result["warnings"],
        )

    def test_six_bullets_is_still_acceptable(self) -> None:
        result = validated(self.build(font_pt=24, bullets=6))
        self.assertEqual(
            [line for line in result["warnings"] if "Many bullets" in line], []
        )

    def test_the_slide_geometry_is_reported(self) -> None:
        result = validated(self.build(font_pt=24, bullets=1))
        # python-pptx's default template is 10" x 7.5", i.e. 4:3.
        self.assertIn('10.0" × 7.5"', "\n".join(result["info"]))


class PdfToImagesTests(TemporaryDirectoryTestCase):
    def test_only_jpeg_and_png_are_accepted(self) -> None:
        for accepted in ("jpg", "jpeg", "png", "PNG"):
            with self.subTest(format=accepted):
                converter = pdf_to_images.PDFToImagesConverter(
                    "deck.pdf", "out/slide", format=accepted
                )
                self.assertEqual(converter.format, accepted.lower())
        for rejected in ("gif", "tiff", "svg"):
            with self.subTest(format=rejected):
                with self.assertRaisesRegex(ValueError, "Unsupported format"):
                    pdf_to_images.PDFToImagesConverter(
                        "deck.pdf", "out/slide", format=rejected
                    )

    def test_a_missing_pdf_is_refused_before_any_conversion(self) -> None:
        converter = pdf_to_images.PDFToImagesConverter(
            str(self.root / "absent.pdf"), str(self.root / "slide")
        )
        with contextlib.redirect_stdout(io.StringIO()):
            with self.assertRaises(FileNotFoundError):
                converter.convert()

    def convert(self, **kwargs):
        if not pdf_to_images.HAS_PYMUPDF:
            self.skipTest("PyMuPDF is not installed")
        source = self.deck(pages=3, size=(1280, 720))
        converter = pdf_to_images.PDFToImagesConverter(
            str(source), str(self.root / "review" / "slide"), **kwargs
        )
        with contextlib.redirect_stdout(io.StringIO()):
            return converter.convert()

    def test_every_page_becomes_one_zero_padded_image(self) -> None:
        produced = self.convert(dpi=72)
        self.assertEqual(
            [path.name for path in produced],
            ["slide-001.jpg", "slide-002.jpg", "slide-003.jpg"],
        )
        for path in produced:
            self.assertTrue(path.is_file())

    def test_the_output_directory_is_created_on_demand(self) -> None:
        produced = self.convert(dpi=72)
        self.assertEqual(produced[0].parent, self.root / "review")

    def test_the_requested_dpi_sets_the_pixel_size(self) -> None:
        # A 720-point-tall slide at 144 dpi is 1440 pixels: zoom = dpi / 72.
        produced = self.convert(dpi=144, format="png")
        with Image.open(produced[0]) as rendered:
            self.assertEqual(rendered.size, (2560, 1440))

    def test_a_page_range_converts_only_those_slides(self) -> None:
        produced = self.convert(dpi=72, first_page=2, last_page=3)
        # 1-indexed and inclusive, and the names keep the original numbering.
        self.assertEqual(
            [path.name for path in produced], ["slide-002.jpg", "slide-003.jpg"]
        )

    def test_the_format_decides_the_encoding_not_just_the_suffix(self) -> None:
        produced = self.convert(dpi=72, format="png")
        with Image.open(produced[0]) as rendered:
            self.assertEqual(rendered.format, "PNG")


class ImageCollectionTests(TemporaryDirectoryTestCase):
    def test_files_are_ordered_by_name_so_numbering_decides_the_deck(self) -> None:
        # Zero-padded names sort correctly; "10" must not precede "2".
        for name in ("03_methods.png", "01_title.png", "10_end.png", "02_intro.png"):
            self.image(name)
        found = slides_to_pdf.get_image_files([str(self.root)])
        self.assertEqual(
            [path.name for path in found],
            ["01_title.png", "02_intro.png", "03_methods.png", "10_end.png"],
        )

    def test_a_directory_yields_every_supported_image(self) -> None:
        for name in ("a.png", "b.jpg", "c.jpeg", "d.bmp"):
            self.image(name)
        found = slides_to_pdf.get_image_files([str(self.root)])
        self.assertEqual(len(found), 4)

    def test_a_non_image_file_is_skipped_with_a_warning(self) -> None:
        self.image("slide.png")
        notes = self.root / "notes.txt"
        notes.write_text("not a slide")
        stream = io.StringIO()
        with contextlib.redirect_stdout(stream):
            found = slides_to_pdf.get_image_files([str(notes), str(self.root / "slide.png")])
        self.assertEqual([path.name for path in found], ["slide.png"])
        self.assertIn("Skipping non-image", stream.getvalue())

    def test_the_same_file_named_twice_appears_once(self) -> None:
        path = self.image("slide.png")
        found = slides_to_pdf.get_image_files([str(path), str(path)])
        self.assertEqual(len(found), 1)

    def test_a_glob_pattern_is_expanded(self) -> None:
        self.image("s1.png")
        self.image("s2.png")
        self.image("other.jpg")
        found = slides_to_pdf.get_image_files([str(self.root / "s*.png")])
        self.assertEqual([path.name for path in found], ["s1.png", "s2.png"])

    def test_a_path_that_matches_nothing_yields_nothing(self) -> None:
        self.assertEqual(slides_to_pdf.get_image_files([str(self.root / "*.png")]), [])


class SlidesToPdfTests(TemporaryDirectoryTestCase):
    def combine(self, images, name="talk.pdf", **kwargs):
        destination = self.root / "out" / name
        with contextlib.redirect_stdout(io.StringIO()) as printed:
            ok = slides_to_pdf.combine_images_to_pdf(images, destination, **kwargs)
        return ok, destination, printed.getvalue()

    def page_count(self, path: Path) -> int:
        if validate_presentation.HAS_PDF_READER:
            from validate_presentation import PdfReader

            return len(PdfReader(str(path)).pages)
        self.skipTest("neither pypdf nor PyPDF2 is installed")

    def test_the_slides_become_pages_in_the_order_given(self) -> None:
        images = [self.image(f"slide{index}.png") for index in range(1, 4)]
        ok, destination, _ = self.combine(images)
        self.assertTrue(ok)
        self.assertEqual(self.page_count(destination), 3)

    def test_the_output_directory_is_created(self) -> None:
        ok, destination, _ = self.combine([self.image("slide.png")])
        self.assertTrue(ok)
        self.assertTrue(destination.is_file())

    def test_no_images_is_refused_rather_than_writing_an_empty_pdf(self) -> None:
        ok, destination, printed = self.combine([])
        self.assertFalse(ok)
        self.assertFalse(destination.exists())
        self.assertIn("No image files found", printed)

    def test_a_transparent_slide_is_flattened_onto_white(self) -> None:
        # PDF has no alpha channel; without the flattening step Pillow raises.
        ok, destination, _ = self.combine([self.image("slide.png", mode="RGBA")])
        self.assertTrue(ok)
        self.assertEqual(self.page_count(destination), 1)

    def test_a_greyscale_slide_is_converted_rather_than_rejected(self) -> None:
        ok, _, _ = self.combine([self.image("grey.png", mode="L")], name="grey.pdf")
        self.assertTrue(ok)

    def test_an_unreadable_slide_aborts_instead_of_producing_a_short_deck(self) -> None:
        broken = self.root / "broken.png"
        broken.write_bytes(b"not an image")
        ok, destination, printed = self.combine([self.image("ok.png"), broken])
        self.assertFalse(ok)
        self.assertFalse(destination.exists())
        self.assertIn("Error loading", printed)


class SlideImageEnvironmentTests(unittest.TestCase):
    """The wrapper hands the credential to the child through the environment."""

    def test_only_allowlisted_variables_reach_the_subprocess(self) -> None:
        environment = {
            "PATH": "/usr/bin",
            "HOME": "/home/someone",
            "AWS_SECRET_ACCESS_KEY": "should-not-be-forwarded",
            "GITHUB_TOKEN": "also-not",
        }
        with mock.patch.dict("os.environ", environment, clear=True):
            built = generate_slide_image.build_subprocess_env(None)
        self.assertEqual(built, {"PATH": "/usr/bin", "HOME": "/home/someone"})

    def test_the_api_key_is_injected_under_the_expected_name(self) -> None:
        with mock.patch.dict("os.environ", {"PATH": "/usr/bin"}, clear=True):
            built = generate_slide_image.build_subprocess_env("sk-or-test")
        self.assertEqual(built["OPENROUTER_API_KEY"], "sk-or-test")

    def test_no_key_means_no_key_variable_rather_than_an_empty_one(self) -> None:
        # An empty value would look configured and fail deep inside the request.
        with mock.patch.dict("os.environ", {"PATH": "/usr/bin"}, clear=True):
            for value in (None, ""):
                with self.subTest(value=value):
                    self.assertNotIn(
                        "OPENROUTER_API_KEY",
                        generate_slide_image.build_subprocess_env(value),
                    )

    def test_the_allowlist_carries_no_credential_variables(self) -> None:
        forwarded = " ".join(generate_slide_image.FORWARDED_ENV_VARS).upper()
        for banned in ("SECRET", "TOKEN", "PASSWORD", "API_KEY", "CREDENTIAL"):
            with self.subTest(term=banned):
                self.assertNotIn(banned, forwarded)

    def test_the_allowlist_matches_the_schematic_wrapper_in_this_skill(self) -> None:
        # Both wrappers in this skill spawn an AI child the same way; a variable
        # added to one and not the other breaks proxies or TLS in half the tools.
        import generate_schematic

        self.assertEqual(
            list(generate_slide_image.FORWARDED_ENV_VARS),
            list(generate_schematic.FORWARDED_ENV_VARS),
        )


class SlideImageInvocationTests(TemporaryDirectoryTestCase):
    """What the wrapper actually spawns, with the subprocess intercepted."""

    def invoke(self, *arguments):
        completed = subprocess.CompletedProcess(args=[], returncode=0)
        argv = ["generate_slide_image.py", *arguments]
        with mock.patch.object(sys, "argv", argv), \
             mock.patch.dict(os.environ, {"OPENROUTER_API_KEY": "sk-or-test"}), \
             mock.patch("subprocess.run", return_value=completed) as run:
            with contextlib.redirect_stdout(io.StringIO()):
                with self.assertRaises(SystemExit) as raised:
                    generate_slide_image.main()
        self.assertEqual(raised.exception.code, 0)
        return run.call_args

    def test_the_generator_script_beside_it_is_what_runs(self) -> None:
        command = self.invoke("a title slide", "-o", str(self.root / "slide.png"))
        self.assertEqual(
            Path(command.args[0][1]).name, "generate_slide_image_ai.py"
        )
        self.assertEqual(command.args[0][0], sys.executable)

    def test_the_credential_never_appears_on_the_command_line(self) -> None:
        command = self.invoke("a slide", "-o", str(self.root / "slide.png"))
        self.assertNotIn("sk-or-test", " ".join(command.args[0]))
        self.assertEqual(command.kwargs["env"]["OPENROUTER_API_KEY"], "sk-or-test")

    def test_the_visual_only_mode_is_forwarded(self) -> None:
        command = self.invoke("a figure", "-o", "f.png", "--visual-only")
        self.assertIn("--visual-only", command.args[0])

    def test_each_attachment_is_forwarded_in_order(self) -> None:
        command = self.invoke(
            "compare these", "-o", "c.png", "--attach", "before.png",
            "--attach", "after.png",
        )
        forwarded = command.args[0]
        self.assertEqual(forwarded.count("--attach"), 2)
        self.assertLess(forwarded.index("before.png"), forwarded.index("after.png"))

    def test_the_iteration_ceiling_is_enforced_by_the_wrapper(self) -> None:
        # Two is the documented maximum; a larger request is clamped, and the
        # default is left implicit rather than passed through.
        self.assertNotIn("--iterations", self.invoke("s", "-o", "s.png").args[0])
        self.assertNotIn(
            "--iterations", self.invoke("s", "-o", "s.png", "--iterations", "5").args[0]
        )
        clamped = self.invoke("s", "-o", "s.png", "--iterations", "1").args[0]
        self.assertEqual(clamped[clamped.index("--iterations") + 1], "1")

    def test_without_a_key_nothing_is_spawned(self) -> None:
        # The wrapper also resolves a credential from any .env file at or above the
        # working directory, so the temporary root stands in for a machine that has
        # none -- otherwise a developer's own .env would satisfy the lookup.
        origin = os.getcwd()
        self.addCleanup(os.chdir, origin)
        os.chdir(self.root)

        argv = ["generate_slide_image.py", "a slide", "-o", "s.png"]
        with mock.patch.object(sys, "argv", argv), \
             mock.patch.dict(os.environ, {}, clear=True), \
             mock.patch("subprocess.run") as run:
            with contextlib.redirect_stdout(io.StringIO()) as printed:
                with self.assertRaises(SystemExit) as raised:
                    generate_slide_image.main()
        self.assertEqual(raised.exception.code, 1)
        self.assertIn("OPENROUTER_API_KEY", printed.getvalue())
        run.assert_not_called()


class SlideGeneratorTestCase(TemporaryDirectoryTestCase):
    """Base class: an offline SlideImageGenerator with a dummy credential."""

    def setUp(self) -> None:
        super().setUp()
        pytest.importorskip("requests", reason="the AI generator needs requests")
        import generate_slide_image_ai

        self.module = generate_slide_image_ai
        self.generator = generate_slide_image_ai.SlideImageGenerator(
            api_key="sk-or-test"
        )

    @staticmethod
    def data_url(payload: bytes) -> str:
        return "data:image/png;base64," + base64.b64encode(payload).decode()


class ApiKeyTests(TemporaryDirectoryTestCase):
    def test_a_generator_without_a_credential_refuses_to_start(self) -> None:
        pytest.importorskip("requests", reason="the AI generator needs requests")
        import generate_slide_image_ai

        # Run from an empty directory: the constructor falls back to a .env in
        # the working directory, and this repository has one.
        origin = os.getcwd()
        self.addCleanup(os.chdir, origin)
        os.chdir(self.root)
        with mock.patch.dict(os.environ, {}, clear=True):
            with self.assertRaisesRegex(ValueError, "OPENROUTER_API_KEY"):
                generate_slide_image_ai.SlideImageGenerator()

    def test_the_environment_variable_is_accepted_in_place_of_an_argument(self) -> None:
        pytest.importorskip("requests", reason="the AI generator needs requests")
        import generate_slide_image_ai

        with mock.patch.dict(os.environ, {"OPENROUTER_API_KEY": "sk-or-env"}):
            generator = generate_slide_image_ai.SlideImageGenerator()
        self.assertEqual(generator.api_key, "sk-or-env")


class GuidelineTests(SlideGeneratorTestCase):
    def test_the_presentation_threshold_is_lower_than_a_paper_figure(self) -> None:
        # Documented as deliberately below the journal/conference bar.
        self.assertEqual(self.module.SlideImageGenerator.QUALITY_THRESHOLD, 6.5)

    def test_a_full_slide_prompt_asks_for_slide_structure(self) -> None:
        guidelines = self.module.SlideImageGenerator.FULL_SLIDE_GUIDELINES
        for requirement in ("16:9", "title", "sans-serif"):
            with self.subTest(requirement=requirement):
                self.assertIn(requirement, guidelines.lower())

    def test_a_visual_only_prompt_does_not_ask_for_slide_furniture(self) -> None:
        # In the PPT workflow the text is added in PowerPoint, so a generated
        # title baked into the image would be duplicated.
        guidelines = self.module.SlideImageGenerator.VISUAL_ONLY_GUIDELINES
        self.assertNotIn("16:9", guidelines)
        self.assertIn("No text unless essential", guidelines)

    def test_the_improved_prompt_carries_the_critique_and_the_right_guidelines(self) -> None:
        for visual_only, expected in (
            (False, self.module.SlideImageGenerator.FULL_SLIDE_GUIDELINES),
            (True, self.module.SlideImageGenerator.VISUAL_ONLY_GUIDELINES),
        ):
            with self.subTest(visual_only=visual_only):
                prompt = self.generator.improve_prompt(
                    "a title slide", "the text is too small", 2, visual_only
                )
                self.assertIn("the text is too small", prompt)
                self.assertIn("a title slide", prompt)
                self.assertIn(expected.strip(), prompt)


class ImageExtractionTests(SlideGeneratorTestCase):
    PAYLOAD = b"\x89PNG\r\n\x1a\n-pretend-this-is-an-image"

    def response(self, message: dict) -> dict:
        return {"choices": [{"message": message}]}

    def test_an_image_in_the_images_field_is_decoded(self) -> None:
        response = self.response(
            {"images": [{"type": "image_url",
                         "image_url": {"url": self.data_url(self.PAYLOAD)}}]}
        )
        self.assertEqual(
            self.generator._extract_image_from_response(response), self.PAYLOAD
        )

    def test_wrapped_base64_is_reassembled(self) -> None:
        # Some responses arrive hard-wrapped; the newlines are not data.
        encoded = base64.b64encode(self.PAYLOAD).decode()
        wrapped = "\n".join(encoded[index:index + 8] for index in range(0, len(encoded), 8))
        response = self.response(
            {"images": [{"type": "image_url",
                         "image_url": {"url": f"data:image/png;base64,{wrapped}"}}]}
        )
        self.assertEqual(
            self.generator._extract_image_from_response(response), self.PAYLOAD
        )

    def test_an_image_embedded_in_the_content_string_is_found(self) -> None:
        response = self.response(
            {"content": f"Here you go: {self.data_url(self.PAYLOAD)} enjoy"}
        )
        self.assertEqual(
            self.generator._extract_image_from_response(response), self.PAYLOAD
        )

    def test_an_image_in_a_content_block_is_found(self) -> None:
        response = self.response(
            {"content": [
                {"type": "text", "text": "here"},
                {"type": "image_url", "image_url": {"url": self.data_url(self.PAYLOAD)}},
            ]}
        )
        self.assertEqual(
            self.generator._extract_image_from_response(response), self.PAYLOAD
        )

    def test_a_response_with_no_image_yields_none(self) -> None:
        for message in ({}, {"content": "I cannot draw that"}, {"images": []}):
            with self.subTest(message=message):
                self.assertIsNone(
                    self.generator._extract_image_from_response(self.response(message))
                )

    def test_a_refusal_with_no_choices_yields_none(self) -> None:
        self.assertIsNone(self.generator._extract_image_from_response({"choices": []}))

    def test_a_truncated_data_url_yields_none_rather_than_raising(self) -> None:
        response = self.response(
            {"images": [{"type": "image_url", "image_url": {"url": "data:image/png"}}]}
        )
        self.assertIsNone(self.generator._extract_image_from_response(response))


class AttachmentEncodingTests(SlideGeneratorTestCase):
    def test_the_mime_type_follows_the_file_extension(self) -> None:
        expected = {
            "chart.png": "image/png",
            "photo.jpg": "image/jpeg",
            "photo.jpeg": "image/jpeg",
            "loop.gif": "image/gif",
            "shot.webp": "image/webp",
        }
        for name, mime in expected.items():
            with self.subTest(name=name):
                path = self.root / name
                path.write_bytes(b"bytes")
                self.assertTrue(
                    self.generator._image_to_base64(str(path)).startswith(
                        f"data:{mime};base64,"
                    )
                )

    def test_an_unknown_extension_falls_back_to_png(self) -> None:
        path = self.root / "figure.tiff"
        path.write_bytes(b"bytes")
        self.assertTrue(
            self.generator._image_to_base64(str(path)).startswith("data:image/png")
        )

    def test_the_encoded_payload_round_trips(self) -> None:
        path = self.root / "chart.png"
        path.write_bytes(b"\x00\x01\x02payload")
        encoded = self.generator._image_to_base64(str(path)).split(",", 1)[1]
        self.assertEqual(base64.b64decode(encoded), b"\x00\x01\x02payload")


class ReviewDecisionTests(SlideGeneratorTestCase):
    def setUp(self) -> None:
        super().setUp()
        self.image_path = self.root / "slide.png"
        self.image_path.write_bytes(b"image bytes")

    def review(self, message: dict):
        with mock.patch.object(
            self.generator, "_make_request", return_value={"choices": [{"message": message}]}
        ):
            return self.generator.review_image(str(self.image_path), "a slide", 1)

    def test_a_score_at_the_threshold_is_accepted(self) -> None:
        # 6.5 is the threshold itself, so it must not trigger another round.
        review = self.review({"content": "SCORE: 6.5\nVERDICT: ACCEPTABLE"})
        self.assertEqual(review.score, 6.5)
        self.assertFalse(review.needs_improvement)

    def test_a_score_below_the_threshold_asks_for_another_round(self) -> None:
        review = self.review({"content": "SCORE: 6.4"})
        self.assertEqual(review.score, 6.4)
        self.assertTrue(review.needs_improvement)

    def test_an_explicit_verdict_overrides_a_high_score(self) -> None:
        # The reviewer's words win: a 9 with NEEDS_IMPROVEMENT is not accepted.
        review = self.review({"content": "SCORE: 9\nVERDICT: NEEDS_IMPROVEMENT"})
        self.assertEqual(review.score, 9.0)
        self.assertTrue(review.needs_improvement)

    def test_a_loosely_worded_score_is_still_read(self) -> None:
        review = self.review({"content": "Overall quality: 8.5 out of 10"})
        self.assertEqual(review.score, 8.5)

    def test_an_unscored_review_does_not_block_the_pipeline(self) -> None:
        # No parseable score and no verdict: still no extra round, because a
        # reviewer that said nothing measurable says nothing about the image.
        # The score stays None -- a stand-in number here would reach the review
        # log indistinguishable from one the reviewer actually gave.
        review = self.review({"content": "Looks good to me"})
        self.assertIsNone(review.score)
        self.assertFalse(review.needs_improvement)
        self.assertTrue(review.error)

    def test_a_review_returned_as_content_blocks_is_flattened(self) -> None:
        review = self.review(
            {"content": [{"type": "text", "text": "SCORE: 6.0"},
                         {"type": "text", "text": "ISSUES: crowded"}]}
        )
        self.assertEqual(review.score, 6.0)
        self.assertIn("crowded", review.critique)

    def test_a_reasoning_only_reply_is_used_as_the_critique(self) -> None:
        review = self.review({"content": "", "reasoning": "SCORE: 6.0"})
        self.assertEqual(review.score, 6.0)
        self.assertIn("SCORE", review.critique)

    def test_a_failed_review_is_not_treated_as_a_bad_image(self) -> None:
        with mock.patch.object(
            self.generator, "_make_request", side_effect=RuntimeError("API down")
        ):
            review = self.generator.review_image(str(self.image_path), "a slide", 1)
        self.assertFalse(review.needs_improvement)
        # Nor as a passing one: the image was never actually measured.
        self.assertIsNone(review.score)
        self.assertFalse(review.reviewed)
        self.assertIn("API down", review.critique)


class RefinementLoopTests(SlideGeneratorTestCase):
    PAYLOAD = b"first-image"
    BETTER = b"second-image"

    def reviewed(self, critique, score, needs_improvement):
        """A completed review, in the shape `review_image` returns."""
        return self.module.ReviewResult(
            critique, score, needs_improvement, reviewed=True
        )

    def generate(self, images, reviews, **kwargs):
        destination = self.root / "slides" / "01_title.png"
        with mock.patch.object(self.generator, "generate_image", side_effect=images) as drew, \
             mock.patch.object(self.generator, "review_image", side_effect=reviews):
            with contextlib.redirect_stdout(io.StringIO()):
                results = self.generator.generate_slide(
                    "a title slide", str(destination), **kwargs
                )
        return results, destination, drew

    def test_a_good_first_image_stops_early_and_is_written(self) -> None:
        results, destination, drew = self.generate(
            [self.PAYLOAD], [self.reviewed("great", 8.0, False)]
        )
        self.assertTrue(results["success"])
        self.assertTrue(results["early_stop"])
        self.assertEqual(results["final_score"], 8.0)
        self.assertEqual(len(results["iterations"]), 1)
        self.assertEqual(drew.call_count, 1)
        # The bytes on disk are the ones that passed review.
        self.assertEqual(destination.read_bytes(), self.PAYLOAD)

    def test_a_weak_image_is_regenerated_from_an_improved_prompt(self) -> None:
        results, destination, drew = self.generate(
            [self.PAYLOAD, self.BETTER],
            [self.reviewed("text too small", 5.0, True), self.reviewed("better", 7.5, False)],
        )
        self.assertTrue(results["success"])
        self.assertEqual(len(results["iterations"]), 2)
        self.assertEqual(destination.read_bytes(), self.BETTER)
        # The second prompt must carry the first critique, or the retry is blind.
        second_prompt = drew.call_args_list[1].args[0]
        self.assertIn("text too small", second_prompt)

    def test_the_iteration_ceiling_still_keeps_the_last_image(self) -> None:
        results, destination, drew = self.generate(
            [self.PAYLOAD, self.BETTER],
            [self.reviewed("bad", 3.0, True), self.reviewed("still bad", 4.0, True)],
            iterations=2,
        )
        self.assertEqual(drew.call_count, 2)
        self.assertFalse(results["early_stop"])
        # Better to hand back the best effort than nothing at all.
        self.assertTrue(results["success"])
        self.assertEqual(destination.read_bytes(), self.BETTER)

    def test_a_failed_generation_writes_no_file_and_reports_failure(self) -> None:
        results, destination, _ = self.generate([None, None], [])
        self.assertFalse(results["success"])
        self.assertIsNone(results["final_image"])
        self.assertFalse(destination.exists())
        self.assertTrue(all(not step["success"] for step in results["iterations"]))

    def test_the_mode_and_threshold_are_recorded_in_the_results(self) -> None:
        results, _, _ = self.generate(
            [self.PAYLOAD], [self.reviewed("fine", 7.0, False)], visual_only=True
        )
        self.assertEqual(results["mode"], "visual_only")
        self.assertEqual(results["quality_threshold"], 6.5)
        self.assertEqual(results["user_prompt"], "a title slide")


class AiCliValidationTests(TemporaryDirectoryTestCase):
    """Argument validation happens before any request is made."""

    def invoke(self, *arguments, key="sk-or-test"):
        environment = child_environment()
        if key is not None:
            environment["OPENROUTER_API_KEY"] = key
        return subprocess.run(
            [sys.executable, str(SCRIPTS / "generate_slide_image_ai.py"), *arguments],
            capture_output=True,
            text=True,
            timeout=120,
            env=environment,
            # The script resolves a credential from any .env file at or above the
            # working directory. Running from the temporary root keeps a developer's
            # real .env out of reach, so key=None genuinely means "no credential".
            cwd=self.root,
        )

    def test_more_iterations_than_allowed_is_refused(self) -> None:
        result = self.invoke(
            "a slide", "-o", str(self.root / "s.png"), "--iterations", "3"
        )
        self.assertEqual(result.returncode, 1)
        self.assertIn("between 1 and 2", result.stdout)
        self.assertFalse((self.root / "s.png").exists())

    def test_zero_iterations_is_refused(self) -> None:
        result = self.invoke(
            "a slide", "-o", str(self.root / "s.png"), "--iterations", "0"
        )
        self.assertEqual(result.returncode, 1)
        self.assertIn("between 1 and 2", result.stdout)

    def test_a_missing_attachment_is_named_before_any_request(self) -> None:
        result = self.invoke(
            "a slide", "-o", str(self.root / "s.png"),
            "--attach", str(self.root / "absent.png"),
        )
        self.assertEqual(result.returncode, 1)
        self.assertIn("Attachment file not found", result.stdout)

    def test_a_missing_credential_is_reported_without_a_traceback(self) -> None:
        result = self.invoke("a slide", "-o", str(self.root / "s.png"), key=None)
        self.assertEqual(result.returncode, 1)
        self.assertIn("OPENROUTER_API_KEY", result.stdout)
        self.assertNotIn("Traceback", result.stderr)


if __name__ == "__main__":
    unittest.main()
